import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
import os

# 1. Считаем данные из Excel
df = pd.read_excel('tasks.xlsx', sheet_name='Таблица_задач')

# Убедимся, что даты в формате datetime
df['Start'] = pd.to_datetime(df['Start'])
df['Finish'] = pd.to_datetime(df['Finish'])

# 2. Построим диаграмму Ганта
fig, ax = plt.subplots(figsize=(10, 6))

tasks = df['Task']
starts = df['Start']
finishes = df['Finish']
durations = finishes - starts

# Для визуализации задач по оси Y
y_pos = range(len(tasks))

# Построение горизонтальных баров
ax.barh(y_pos, durations.dt.days, left=starts.map(mdates.date2num), height=0.5, align='center', color='skyblue')

# Настройка оси Y
ax.set_yticks(y_pos)
ax.set_yticklabels(tasks)
ax.invert_yaxis()  # Чтобы задачи шли сверху вниз

# Форматирование оси X (даты)
ax.xaxis_date()
ax.xaxis.set_major_locator(mdates.AutoDateLocator())
ax.xaxis.set_major_formatter(mdates.DateFormatter('%Y-%m-%d'))
plt.xticks(rotation=45)

# Добавим подписи
for i, (start, duration) in enumerate(zip(starts, durations)):
    ax.text(mdates.date2num(start) + duration.days / 2, i, f'{duration.days} дн.', va='center', ha='center', color='black')

plt.title('Диаграмма Ганта')
plt.tight_layout()

# 3. Сохраним диаграмму в изображение
img_path = 'gantt_chart.png'
plt.savefig(img_path, dpi=150)
plt.close()

# 4. Создадим презентацию и вставим изображение
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Пустой слайд

left = Inches(1)
top = Inches(1)
height = Inches(5)

slide.shapes.add_picture(img_path, left, top, height=height)

# 5. Сохраним файл как tasks.ppsx
prs.save('tasks.pptx')

# Очистим временный файл изображения
os.remove(img_path)

print("Презентация tasks.pptx успешно создана.")
