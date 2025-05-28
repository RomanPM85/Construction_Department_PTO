from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
import numpy as np
import matplotlib.pyplot as plt

# Функция для создания графика и его сохранения как изображения
def create_chart():
    # Пример данных
    x = np.linspace(0, 10, 100)
    y = np.sin(x)

    plt.figure()
    plt.plot(x, y, label='sin(x)')
    plt.title('График функции')
    plt.xlabel('x')
    plt.ylabel('sin(x)')
    plt.legend()
    plt.grid()
    plt.savefig('chart.png')  # Сохраняем график как изображение
    plt.close()

# Создание новой презентации
presentation = Presentation()

# Слайд 1: Заголовок
slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])  # Используем макет с заголовком
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Пример презентации"
subtitle.text = "Создано с помощью python-pptx"

# Слайд 2: Изображение
slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])  # Пустой макет
img_path = 'example_image.jpg'  # Замените на путь к вашему изображению
slide2.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(5))

# Слайд 3: Таблица
slide3 = presentation.slides.add_slide(presentation.slide_layouts[5])  # Пустой макет
table = slide3.shapes.add_table(rows=3, cols=2, left=Inches(1), top=Inches(1), width=Inches(5), height=Inches(2)).table

# Заполнение таблицы
table.cell(0, 0).text = 'Заголовок 1'
table.cell(0, 1).text = 'Заголовок 2'
table.cell(1, 0).text = 'Данные 1'
table.cell(1, 1).text = 'Данные 2'
table.cell(2, 0).text = 'Данные 3'
table.cell(2, 1).text = 'Данные 4'

# Слайд 4: Текст
slide4 = presentation.slides.add_slide(presentation.slide_layouts[1])  # Макет с заголовком и текстом
title = slide4.shapes.title
content = slide4.placeholders[1]

title.text = "Текстовый слайд"
content.text = "Это пример текстового слайда, созданного с помощью python-pptx."

# Слайд 5: График
create_chart()  # Создаем график
slide5 = presentation.slides.add_slide(presentation.slide_layouts[5])  # Пустой макет
slide5.shapes.add_picture('chart.png', Inches(1), Inches(1), width=Inches(5))

# Сохранение презентации
presentation.save('example_presentation.pptx')
print("Презентация сохранена как example_presentation.pptx")
