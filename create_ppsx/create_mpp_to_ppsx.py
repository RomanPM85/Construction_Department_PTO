import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.dates import date2num
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
import io

# Загрузка данных из CSV (экспортированного из .mpp)
df = pd.read_csv('tasks.csv', parse_dates=['Start', 'Finish'])

# Создание графика Ганта
fig, ax = plt.subplots(figsize=(10, 6))

for i, task in enumerate(df.itertuples()):
    start_num = date2num(task.Start)
    finish_num = date2num(task.Finish)
    ax.barh(i, finish_num - start_num, left=start_num, height=0.4, align='center')
    ax.text(start_num, i, task.Task, va='center', ha='right', fontsize=8)

ax.set_yticks(range(len(df)))
ax.set_yticklabels(df['Task'])
ax.invert_yaxis()
ax.xaxis_date()
plt.xlabel('Дата')
plt.title('График Ганта')

# Сохранение графика в буфер
img_stream = io.BytesIO()
plt.savefig(img_stream, format='png', bbox_inches='tight')
plt.close(fig)
img_stream.seek(0)

# Создание презентации PowerPoint
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Пустой слайд

# Вставка изображения графика
left = Inches(1)
top = Inches(1)
height = Inches(4.5)
slide.shapes.add_picture(img_stream, left, top, height=height)

# Сохранение в .ppsx (презентация в режиме показа)
prs.save('Gantt_chart.ppsx')

print("Презентация с графиком Ганта сохранена в Gantt_chart.ppsx")
